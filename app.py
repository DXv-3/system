#!/usr/bin/env python3
"""
LLM Pre-Processing Pipeline Generator

A Flask web application that converts uploaded ZIP files into LLM-optimized
packages by applying the LLM Data-Type Hierarchy to each file.

DEPENDENCIES:
-------------
System packages (install via):
    sudo apt-get install pandoc tesseract-ocr poppler-utils

Python packages (install via pip):
    pip install Flask python-docx python-pptx openpyxl Pillow pytesseract

USAGE:
------
    python app.py

Then navigate to http://localhost:5000 in your browser.
"""

import os
import sys
import uuid
import zipfile
import subprocess
import shutil
import logging
from datetime import datetime
from pathlib import Path

from flask import Flask, request, send_file, jsonify, render_template_string

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# ============================================================================
# Configuration
# ============================================================================
UPLOAD_FOLDER = Path("/tmp/llm_pipeline_uploads")
EXTRACT_FOLDER = Path("/tmp/llm_pipeline_extract")
OUTPUT_FOLDER = Path("/tmp/llm_pipeline_output")

# Ensure directories exist
for folder in [UPLOAD_FOLDER, EXTRACT_FOLDER, OUTPUT_FOLDER]:
    folder.mkdir(parents=True, exist_ok=True)

# Check system dependencies
def check_system_dependencies():
    """Check that required system tools are available."""
    missing = []
    for cmd, name in [('pandoc', 'pandoc'), ('tesseract', 'Tesseract OCR'), ('pdftotext', 'poppler-utils')]:
        try:
            subprocess.run(['which', cmd], capture_output=True, check=True)
        except (subprocess.CalledProcessError, FileNotFoundError):
            missing.append(f'{name} ({cmd})')
    return missing

# ============================================================================
# File Type Definitions
# ============================================================================

# Tier S: Raw text passthrough - copy as-is
TIER_S_EXTENSIONS = {
    # Source code
    '.py', '.js', '.ts', '.jsx', '.tsx', '.java', '.go', '.rs', '.c', '.cpp',
    '.h', '.hpp', '.cs', '.php', '.rb', '.swift', '.kt', '.scala', '.r', '.m',
    '.pl', '.sh', '.bash', '.zsh', '.fish', '.ps1', '.bat', '.cmd',
    # Text files
    '.txt', '.text',
    # Data formats
    '.json', '.xml', '.csv', '.tsv', '.yaml', '.yml', '.toml', '.ini',
    '.cfg', '.conf', '.properties',
    # Markup and styling
    '.html', '.htm', '.css', '.scss', '.sass', '.less', '.md', '.markdown',
    '.rst', '.adoc', '.asciidoc',
    # SQL and query languages
    '.sql', '.graphql', '.gql',
    # Config files
    '.gitignore', '.dockerignore', '.env', '.editorconfig', '.prettierrc',
    '.eslintrc', '.babelrc',
}

# Tier A: Convert to Markdown (structured documents)
TIER_A_EXTENSIONS = {
    '.docx', '.doc', '.rtf', '.odt',  # Word processors
    '.pptx', '.ppt', '.odp',         # Presentations
    '.xlsx', '.xls', '.ods',         # Spreadsheets
}

# Tier C: OCR to text (visual extraction)
TIER_C_EXTENSIONS = {
    '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif', '.webp',
}

# Tier F: Ignore/binary - excluded
TIER_F_EXCLUDED = {
    # Executables
    '.exe', '.msi', '.bin', '.run', '.app', '.bat', '.cmd', '.sh',
    # Archives (shouldn't appear but just in case)
    '.zip', '.tar', '.gz', '.bz2', '.7z', '.rar',
    # Media files
    '.mp4', '.avi', '.mov', '.wmv', '.flv', '.webm', '.mkv',
    '.mp3', '.wav', '.flac', '.aac', '.ogg', '.wma',
    # Video/audio formats
    '.mpeg', '.mpg', '.m4a', '.m4v', '.3gp',
    # Fonts
    '.ttf', '.otf', '.woff', '.woff2', '.eot',
    # Database files
    '.db', '.sqlite', '.sqlite3', '.mdb', '.accdb',
    # Virtual machine / system files
    '.iso', '.vmdk', '.ova', '.qcow2',
    # Other binaries
    '.pdf',  # PDFs are handled separately (could be image-based)
}

# ============================================================================
# Helper Functions
# ============================================================================

def get_file_extension(filename: str) -> str:
    """Extract and normalize file extension to lowercase."""
    return Path(filename).suffix.lower()

def should_process_file(filename: str) -> bool:
    """Check if a file should be processed or ignored."""
    ext = get_file_extension(filename)
    if ext in TIER_F_EXCLUDED:
        return False
    return True

def get_tier(filename: str) -> str:
    """Determine which tier a file belongs to."""
    ext = get_file_extension(filename)
    if ext in TIER_S_EXTENSIONS:
        return 'S'
    elif ext in TIER_A_EXTENSIONS:
        return 'A'
    elif ext in TIER_C_EXTENSIONS:
        return 'C'
    elif ext == '.pdf':
        # PDFs are special - try text extraction first, fallback to OCR
        return 'C'
    else:
        # Unknown format - treat as binary/ignore
        return 'F'

def safe_filename(filename: str) -> str:
    """Sanitize filename for safe filesystem usage."""
    # Remove path traversal attempts
    name = Path(filename).name
    # Remove leading dots and slashes
    name = name.lstrip('.').lstrip('/').lstrip('\\')
    return name

def create_unique_dir(base_path: Path) -> Path:
    """Create a unique temporary directory for processing."""
    uid = str(uuid.uuid4())[:8]
    dir_path = base_path / uid
    dir_path.mkdir(parents=True, exist_ok=True)
    return dir_path

# ============================================================================
# Tier Handlers
# ============================================================================

def handle_tier_s(input_path: Path, output_path: Path, rel_path: Path, 
                  manifest_log: list) -> bool:
    """Tier S handler: Copy file as-is (raw text passthrough)."""
    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(input_path, output_path)
        manifest_log.append({
            'original': str(rel_path),
            'output': str(rel_path),
            'action': 'Copied as-is (Tier S)',
            'status': 'success'
        })
        logger.info(f"Tier S: Copied {rel_path}")
        return True
    except Exception as e:
        manifest_log.append({
            'original': str(rel_path),
            'output': '',
            'action': f'Failed to copy: {str(e)}',
            'status': 'error'
        })
        logger.error(f"Tier S error for {rel_path}: {e}")
        return False

def handle_tier_a(input_path: Path, output_path: Path, rel_path: Path,
                  manifest_log: list) -> bool:
    """Tier A handler: Convert document to Markdown using pandoc."""
    try:
        output_path = output_path.with_suffix('.md')
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Use pandoc for conversion
        # Pandoc handles .docx, .rtf, .odt, .pptx, .xlsx, etc.
        result = subprocess.run(
            [
                'pandoc',
                str(input_path),
                '-t', 'markdown',
                '-o', str(output_path),
                '--wrap=preserve',
                '--standalone'
            ],
            capture_output=True,
            text=True,
            timeout=30
        )

        output_rel_path = rel_path.with_suffix('.md')

        if result.returncode == 0 and output_path.exists():
            # If output is empty, try alternative approach
            if output_path.stat().st_size == 0:
                # For Office files, use Python libraries as fallback
                if input_path.suffix.lower() in ['.docx', '.pptx', '.xlsx']:
                    success = _handle_tier_a_python(input_path, output_path)
                    if success:
                        manifest_log.append({
                            'original': str(rel_path),
                            'output': str(output_rel_path),
                            'action': 'Converted to Markdown (Tier A - Python fallback)',
                            'status': 'success'
                        })
                    else:
                        manifest_log.append({
                            'original': str(rel_path),
                            'output': '',
                            'action': 'Conversion failed (Tier A)',
                            'status': 'error'
                        })
                    return success

            manifest_log.append({
                'original': str(rel_path),
                'output': str(output_rel_path),
                'action': 'Converted to Markdown (Tier A)',
                'status': 'success'
            })
            logger.info(f"Tier A: Converted {rel_path}")
            return True
        else:
            # Pandoc failed, try Python libraries
            logger.warning(f"Pandoc failed for {rel_path}: {result.stderr}")
            success = _handle_tier_a_python(input_path, output_path)
            if success:
                manifest_log.append({
                    'original': str(rel_path),
                    'output': str(output_rel_path),
                    'action': 'Converted to Markdown (Tier A - Python fallback)',
                    'status': 'success'
                })
            else:
                manifest_log.append({
                    'original': str(rel_path),
                    'output': '',
                    'action': 'Conversion failed (Tier A)',
                    'status': 'error'
                })
            return success

    except subprocess.TimeoutExpired:
        logger.error(f"Pandoc timeout for {rel_path}")
        manifest_log.append({
            'original': str(rel_path),
            'output': '',
            'action': 'Conversion timeout (Tier A)',
            'status': 'error'
        })
        return False
    except Exception as e:
        logger.error(f"Tier A error for {rel_path}: {e}")
        manifest_log.append({
            'original': str(rel_path),
            'output': '',
            'action': f'Conversion error: {str(e)}',
            'status': 'error'
        })
        return False

def _handle_tier_a_python(input_path: Path, output_path: Path) -> bool:
    """Python-based fallback for Tier A conversions."""
    try:
        ext = input_path.suffix.lower()
        text_content = ""

        if ext == '.docx':
            try:
                import docx
                doc = docx.Document(input_path)
                for para in doc.paragraphs:
                    text_content += para.text + "\n"
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            text_content += cell.text + "\t"
                        text_content += "\n"
            except ImportError:
                return False

        elif ext in ['.pptx', '.ppt']:
            try:
                import pptx
                prs = pptx.Presentation(input_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content += shape.text + "\n"
                    text_content += "\n---\n\n"
            except ImportError:
                return False

        elif ext in ['.xlsx', '.xls']:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(input_path, data_only=True)
                for sheet in wb.worksheets:
                    text_content += f"Sheet: {sheet.title}\n"
                    for row in sheet.iter_rows(values_only=True):
                        row_text = "\t".join([str(cell) if cell is not None else '' for cell in row])
                        text_content += row_text + "\n"
                    text_content += "\n"
            except ImportError:
                return False

        else:
            return False

        if text_content:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text_content)
            return True
        return False

    except Exception as e:
        logger.error(f"Python fallback error for {input_path}: {e}")
        return False

def handle_tier_c(input_path: Path, output_path: Path, rel_path: Path,
                  manifest_log: list) -> bool:
    """Tier C handler: Extract text via OCR (images, PDFs)."""
    try:
        text_content = ""
        ext = input_path.suffix.lower()

        if ext == '.pdf':
            # Try pdftotext first (for text-based PDFs)
            try:
                result = subprocess.run(
                    ['pdftotext', str(input_path), '-'],
                    capture_output=True,
                    text=True,
                    timeout=30
                )
                if result.returncode == 0 and result.stdout.strip():
                    text_content = result.stdout
                    logger.info(f"PDF text extracted via pdftotext: {rel_path}")
                else:
                    # PDF likely image-based, convert to images and OCR
                    text_content = _pdf_ocr(input_path)
            except (subprocess.SubprocessError, FileNotFoundError):
                # pdftotext not available, use OCR directly
                text_content = _pdf_ocr(input_path)

        else:
            # Image file - use pytesseract
            try:
                from PIL import Image
                import pytesseract
                img = Image.open(input_path)
                text_content = pytesseract.image_to_string(img)
            except ImportError:
                logger.error("Pillow or pytesseract not installed")
                manifest_log.append({
                    'original': str(rel_path),
                    'output': '',
                    'action': 'OCR failed: missing dependencies',
                    'status': 'error'
                })
                return False
            except Exception as e:
                logger.error(f"OCR error for {rel_path}: {e}")
                text_content = ""

        output_path = output_path.with_suffix(output_path.suffix + '.txt')
        output_path.parent.mkdir(parents=True, exist_ok=True)

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text_content)

        output_rel_path = rel_path.with_suffix(rel_path.suffix + '.txt')
        manifest_log.append({
            'original': str(rel_path),
            'output': str(output_rel_path),
            'action': 'Text extracted via OCR (Tier C)',
            'status': 'success'
        })
        logger.info(f"Tier C: OCR extracted {rel_path}")
        return True

    except Exception as e:
        logger.error(f"Tier C error for {rel_path}: {e}")
        manifest_log.append({
            'original': str(rel_path),
            'output': '',
            'action': f'OCR error: {str(e)}',
            'status': 'error'
        })
        return False

def _pdf_ocr(input_path: Path) -> str:
    """Convert PDF to images and perform OCR."""
    try:
        from PIL import Image
        import pytesseract
        import tempfile

        text_content = ""

        # Convert PDF to images using pdftoppm or similar
        with tempfile.TemporaryDirectory() as tmpdir:
                result = subprocess.run(
                    ['pdftoppm', '-png', str(input_path),
                     os.path.join(tmpdir, 'page')],
                    capture_output=True
                )
                if result.returncode == 0:
                    for img_file in sorted(Path(tmpdir).glob('page-*.png')):
                        img = Image.open(img_file)
                        text_content += pytesseract.image_to_string(img) + "\n\n"
                else:
                    logger.warning("pdftoppm not available, PDF OCR skipped")

        return text_content

    except Exception as e:
        logger.error(f"PDF OCR error: {e}")
        return ""

def handle_tier_f(filename: str, manifest_log: list) -> None:
    """Tier F handler: Log ignored binary files."""
    manifest_log.append({
        'original': filename,
        'output': '',
        'action': 'Ignored (Binary/Unsupported)',
        'status': 'skipped'
    })
    logger.info(f"Tier F: Ignored {filename}")

# ============================================================================
# Main Processing Pipeline
# ============================================================================

def process_zip(zip_path: Path, output_dir: Path):
    """Main processing pipeline for a ZIP file."""
    manifest_log = []
    zip_name = zip_path.name
    processing_dir = create_unique_dir(EXTRACT_FOLDER)
    output_temp = create_unique_dir(OUTPUT_FOLDER)

    try:
        # Extract ZIP
        logger.info(f"Extracting {zip_name}")
        with zipfile.ZipFile(zip_path, 'r') as zf:
            zf.extractall(processing_dir)

        # Process each file
        total_files = 0
        for root, dirs, files in os.walk(processing_dir):
            for filename in files:
                total_files += 1
                input_path = Path(root) / filename
                rel_path = input_path.relative_to(processing_dir)

                # Skip hidden/system files
                if filename.startswith('.') or filename.startswith('~'):
                    continue

                if not should_process_file(filename):
                    handle_tier_f(filename, manifest_log)
                    continue

                tier = get_tier(filename)
                output_path = output_temp / rel_path

                if tier == 'S':
                    handle_tier_s(input_path, output_path, rel_path, manifest_log)
                elif tier == 'A':
                    handle_tier_a(input_path, output_path, rel_path, manifest_log)
                elif tier == 'C':
                    handle_tier_c(input_path, output_path, rel_path, manifest_log)
                else:
                    handle_tier_f(filename, manifest_log)

        # Generate manifest
        logger.info("Generating manifest")
        manifest_path = output_temp / '_manifest.md'
        generate_manifest(manifest_path, zip_name, output_temp, manifest_log, total_files)

        # Create output ZIP
        output_zip = output_dir / f"llm_optimized_{Path(zip_path).stem}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"
        with zipfile.ZipFile(output_zip, 'w', zipfile.ZIP_DEFLATED) as zf:
            for root, dirs, files in os.walk(output_temp):
                for filename in files:
                    file_path = Path(root) / filename
                    arcname = file_path.relative_to(output_temp)
                    zf.write(file_path, arcname)

        logger.info(f"Output ZIP created: {output_zip}")
        return output_zip

    finally:
        # Cleanup
        logger.info("Cleaning up temporary directories")
        for temp_dir in [processing_dir, output_temp]:
            try:
                if temp_dir.exists():
                    shutil.rmtree(temp_dir)
            except Exception as e:
                logger.warning(f"Failed to cleanup {temp_dir}: {e}")


def generate_manifest(manifest_path: Path, original_zip: str, output_dir: Path,
                      manifest_log: list, total_files: int):
    """Generate the _manifest.md file."""
    lines = []
    lines.append("# LLM Processing Manifest\n")
    lines.append(f"**Original ZIP:** `{original_zip}`\n")
    lines.append(f"**Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
    lines.append(f"**Total Files Processed:** {total_files}\n")
    lines.append(f"**Total Output Files:** {sum(1 for entry in manifest_log if entry['status'] in ('success', 'skipped'))}\n")
    lines.append("---\n\n")

    lines.append("## File Tree\n\n")
    lines.append("```\n")
    lines.append(generate_file_tree(output_dir))
    lines.append("```\n\n")

    lines.append("## Processing Log\n\n")
    lines.append("| # | Original File | Output File | Action | Status |\n")
    lines.append("|---|--------------|------------|--------|--------|\n")
    for i, entry in enumerate(manifest_log, 1):
        lines.append(
            f"| {i} | `{entry['original']}` | "
            f"`{entry['output'] if entry['output'] else '-'}` | "
            f"{entry['action']} | {entry['status']} |\n"
        )

    lines.append("\n---\n")
    lines.append("\n### Summary\n\n")
    
    status_counts = {}
    for entry in manifest_log:
        status = entry['status']
        status_counts[status] = status_counts.get(status, 0) + 1

    for status, count in status_counts.items():
        lines.append(f"- **{status.title()}:** {count}\n")

    lines.append("\n---\n")
    lines.append("\n*Generated by LLM Pre-Processing Pipeline*\n")

    with open(manifest_path, 'w', encoding='utf-8') as f:
        f.writelines(lines)


def generate_file_tree(directory: Path) -> str:
    """Generate a visual file tree representation."""
    tree_lines = []

    def add_tree(prefix: str, path: Path, is_last: bool):
        connector = "└── " if is_last else "├── "
        name = path.name
        tree_lines.append(f"{prefix}{connector}{name}")

        if path.is_dir():
            children = sorted([c for c in path.iterdir()], key=lambda x: (not x.is_dir(), x.name.lower()))
            extension = "    " if is_last else "│   "
            for i, child in enumerate(children):
                add_tree(prefix + extension, child, i == len(children) - 1)
    children = sorted([c for c in directory.iterdir()], key=lambda x: (not x.is_dir(), x.name.lower()))
    for i, child in enumerate(children):
        add_tree("", child, i == len(children) - 1)

    return "\n".join(tree_lines) + "\n"


# ============================================================================
# Flask Application
# ============================================================================

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 500 * 1024 * 1024  # 500MB max upload

HTML_TEMPLATE = """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>LLM Pre-Processing Pipeline</title>
    <style>
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        body {
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Oxygen, Ubuntu, sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
            display: flex;
            justify-content: center;
            align-items: center;
            padding: 20px;
        }
        .container {
            background: white;
            border-radius: 20px;
            padding: 40px;
            max-width: 800px;
            width: 100%;
            box-shadow: 0 20px 60px rgba(0,0,0,0.3);
            animation: slideUp 0.5s ease;
        }
        @keyframes slideUp {
            from {
                opacity: 0;
                transform: translateY(30px);
            }
            to {
                opacity: 1;
                transform: translateY(0);
            }
        }
        h1 {
            color: #2d3748;
            margin-bottom: 10px;
            font-size: 28px;
        }
        .subtitle {
            color: #718096;
            margin-bottom: 30px;
            font-size: 16px;
        }
        .drop-zone {
            border: 3px dashed #cbd5e0;
            border-radius: 15px;
            padding: 60px 40px;
            text-align: center;
            transition: all 0.3s ease;
            cursor: pointer;
            background: #f8fafc;
            position: relative;
            overflow: hidden;
        }
        .drop-zone:hover {
            border-color: #667eea;
            background: #f0f4ff;
        }
        .drop-zone.dragover {
            border-color: #667eea;
            background: #e6eeff;
            transform: scale(1.02);
        }
        .drop-zone .icon {
            font-size: 60px;
            margin-bottom: 15px;
        }
        .drop-zone h3 {
            color: #2d3748;
            margin-bottom: 10px;
            font-size: 20px;
        }
        .drop-zone p {
            color: #718096;
            font-size: 14px;
            line-height: 1.6;
        }
        .file-info {
            display: none;
            margin-top: 20px;
            padding: 15px;
            background: #edf2f7;
            border-radius: 10px;
            color: #2d3748;
        }
        .btn {
            display: none;
            width: 100%;
            padding: 16px;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            border: none;
            border-radius: 10px;
            font-size: 18px;
            font-weight: 600;
            cursor: pointer;
            transition: all 0.3s ease;
            margin-top: 20px;
        }
        .btn:hover {
            transform: translateY(-2px);
            box-shadow: 0 10px 30px rgba(102, 126, 234, 0.4);
        }
        .btn:active {
            transform: translateY(0);
        }
        .btn:disabled {
            opacity: 0.6;
            cursor: not-allowed;
            transform: none;
        }
        .progress-bar {
            display: none;
            width: 100%;
            height: 8px;
            background: #e2e8f0;
            border-radius: 4px;
            margin-top: 20px;
            overflow: hidden;
        }
        .progress-fill {
            height: 100%;
            background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
            width: 0%;
            transition: width 0.3s ease;
        }
        .status-message {
            display: none;
            margin-top: 20px;
            padding: 15px;
            border-radius: 10px;
            text-align: center;
            font-weight: 500;
        }
        .status-message.success {
            background: #c6f6d5;
            color: #22543d;
        }
        .status-message.error {
            background: #fed7d7;
            color: #742a2a;
        }
        .spinner {
            display: none;
            width: 40px;
            height: 40px;
            border: 4px solid #e2e8f0;
            border-top-color: #667eea;
            border-radius: 50%;
            animation: spin 1s linear infinite;
            margin: 20px auto;
        }
        @keyframes spin {
            to { transform: rotate(360deg); }
        }
        .tiers-info {
            margin-top: 30px;
            padding-top: 20px;
            border-top: 1px solid #e2e8f0;
        }
        .tiers-info h4 {
            color: #2d3748;
            margin-bottom: 15px;
            font-size: 16px;
        }
        .tier-badge {
            display: inline-block;
            padding: 4px 12px;
            border-radius: 20px;
            font-size: 12px;
            font-weight: 600;
            margin: 4px;
        }
        .tier-s { background: #c6f6d5; color: #22543d; }
        .tier-a { background: #bee3f8; color: #2c5282; }
        .tier-c { background: #feebc8; color: #7c4d00; }
        .tier-f { background: #e9e9e9; color: #718096; }
        .input-file {
            display: none;
        }
    </style>
</head>
<body>
    <div class="container">
        <h1>LLM Pre-Processing Pipeline</h1>
        <p class="subtitle">Upload a ZIP file and automatically convert its contents for optimal LLM ingestion</p>

        <form id="uploadForm" enctype="multipart/form-data">
            <label class="drop-zone" id="dropZone">
                <div class="icon">📁</div>
                <h3>Click or drag ZIP file here</h3>
                <p>Drop your ZIP archive and we'll process each file according to the LLM Data-Type Hierarchy</p>
                <input type="file" name="file" id="fileInput" class="input-file" accept=".zip" required>
            </label>

            <div class="file-info" id="fileInfo"></div>
            <div class="spinner" id="spinner"></div>
            <div class="progress-bar" id="progressBar">
                <div class="progress-fill" id="progressFill"></div>
            </div>

            <button type="submit" class="btn" id="uploadBtn">Process ZIP File</button>
        </form>

        <div class="status-message" id="statusMessage"></div>

        <div class="tiers-info">
            <h4>Processing Tiers</h4>
            <span class="tier-badge tier-s">Tier S - Raw Text (Passthrough)</span>
            <span class="tier-badge tier-a">Tier A → Markdown</span>
            <span class="tier-badge tier-c">Tier C → OCR Text</span>
            <span class="tier-badge tier-f">Tier F - Ignored</span>
        </div>
    </div>

    <script>
        const dropZone = document.getElementById('dropZone');
        const fileInput = document.getElementById('fileInput');
        const fileInfo = document.getElementById('fileInfo');
        const uploadBtn = document.getElementById('uploadBtn');
        const form = document.getElementById('uploadForm');
        const spinner = document.getElementById('spinner');
        const progressBar = document.getElementById('progressBar');
        const progressFill = document.getElementById('progressFill');
        const statusMessage = document.getElementById('statusMessage');

         let selectedFile = null;
        let depsChecked = false;

        // Check dependencies on load
        checkDependencies();

        function checkDependencies() {
            fetch('/check-deps')
                .then(r => r.json())
                .then(data => {
                    if (data.status !== 'ok') {
                        showStatus(
                            'Missing dependencies: ' + data.missing.join(', ') + 
                            '. ' + data.install_cmd,
                            'error'
                        );
                        uploadBtn.disabled = true;
                    }
                    depsChecked = true;
                })
                .catch(() => {
                    showStatus('Could not check dependencies', 'error');
                    depsChecked = true;
                });
        }


        // Click to select
        dropZone.addEventListener('click', () => fileInput.click());

        // File selection
        fileInput.addEventListener('change', (e) => {
            handleFileSelection(e.target.files[0]);
        });

        // Drag and drop
        dropZone.addEventListener('dragover', (e) => {
            e.preventDefault();
            dropZone.classList.add('dragover');
        });

        dropZone.addEventListener('dragleave', () => {
            dropZone.classList.remove('dragover');
        });

        dropZone.addEventListener('drop', (e) => {
            e.preventDefault();
            dropZone.classList.remove('dragover');
            const files = e.dataTransfer.files;
            if (files.length > 0) {
                handleFileSelection(files[0]);
            }
        });

        function handleFileSelection(file) {
            if (!file) return;

            if (!file.name.endsWith('.zip')) {
                showStatus('Please upload a ZIP file', 'error');
                return;
            }

            selectedFile = file;
            fileInfo.style.display = 'block';
            fileInfo.innerHTML = `<strong>Selected:</strong> ${file.name} (${(file.size / 1024 / 1024).toFixed(2)} MB)`;
            uploadBtn.style.display = 'block';
        }

        // Form submission
        form.addEventListener('submit', async (e) => {
            e.preventDefault();
            if (!selectedFile) return;

            uploadBtn.disabled = true;
            spinner.style.display = 'block';
            progressBar.style.display = 'block';
            statusMessage.style.display = 'none';

            const formData = new FormData();
            formData.append('file', selectedFile);

            try {
                const xhr = new XMLHttpRequest();

                xhr.upload.addEventListener('progress', (e) => {
                    if (e.lengthComputable) {
                        const percent = Math.round((e.loaded / e.total) * 100);
                        progressFill.style.width = percent + '%';
                    }
                });

                xhr.addEventListener('load', () => {
                    spinner.style.display = 'none';

                    if (xhr.status === 200) {
                        const response = JSON.parse(xhr.responseText);

                        // Create download link
                        const link = document.createElement('a');
                        link.href = '/download/' + response.filename;
                        link.download = response.filename;
                        link.click();

                        showStatus('Processing complete! Download started...', 'success');
                        resetForm();

                        // Also fetch the manifest
                        setTimeout(() => {
                            fetch('/download/' + response.filename.replace('.zip', '/_manifest.md'))
                                .then(r => r.text())
                                .then(text => console.log('Manifest:', text))
                                .catch(() => {});
                        }, 1000);
                    } else {
                        const error = JSON.parse(xhr.responseText);
                        showStatus('Error: ' + error.error, 'error');
                        uploadBtn.disabled = false;
                    }
                });

                xhr.addEventListener('error', () => {
                    spinner.style.display = 'none';
                    showStatus('Network error occurred', 'error');
                    uploadBtn.disabled = false;
                });

                xhr.open('POST', '/process');
                xhr.send(formData);

            } catch (error) {
                spinner.style.display = 'none';
                showStatus('Error: ' + error.message, 'error');
                uploadBtn.disabled = false;
            }
        });

        function showStatus(message, type) {
            statusMessage.textContent = message;
            statusMessage.className = 'status-message ' + type;
            statusMessage.style.display = 'block';
        }

        function resetForm() {
            form.reset();
            selectedFile = null;
            fileInfo.style.display = 'none';
            uploadBtn.style.display = 'none';
            progressFill.style.width = '0%';
            setTimeout(() => {
                progressBar.style.display = 'none';
            }, 500);
        }
    </script>
</body>
</html>
"""

@app.route('/')
def index():
    """Serve the main page."""
    return render_template_string(HTML_TEMPLATE)

@app.route('/check-deps')
def check_deps():
    """Check if all system dependencies are installed."""
    missing = check_system_dependencies()
    if missing:
        return jsonify({
            'status': 'missing',
            'missing': missing,
            'install_cmd': 'sudo apt-get install -y pandoc tesseract-ocr poppler-utils'
        })
    return jsonify({'status': 'ok', 'missing': []})

@app.route('/process', methods=['POST'])
def process():
    """Handle ZIP file upload and processing."""
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400

    if not file.filename.lower().endswith('.zip'):
        return jsonify({'error': 'File must be a ZIP archive'}), 400

    # Check system dependencies
    missing_deps = check_system_dependencies()
    if missing_deps:
        return jsonify({
            'error': 'Missing system dependencies',
            'details': 'Required tools not found: ' + ', '.join(missing_deps),
            'install': 'Run: sudo apt-get install -y pandoc tesseract-ocr poppler-utils'
        }), 503

    try:
        # Check file size (soft limit: 500MB)
        file.seek(0, 2)  # Seek to end
        file_size = file.tell()
        file.seek(0)  # Seek back to start
        
        if file_size > 500 * 1024 * 1024:
            return jsonify({
                'error': 'File too large',
                'max_size': '500MB',
                'actual_size': f'{file_size / (1024*1024):.2f}MB'
            }), 413

        # Save uploaded file
        upload_path = UPLOAD_FOLDER / file.filename
        file.save(upload_path)
        logger.info(f"Received file: {upload_path}")

        # Process the ZIP
        output_zip = process_zip(upload_path, OUTPUT_FOLDER)

        # Clean up uploaded file
        try:
            upload_path.unlink()
        except Exception as e:
            logger.warning(f"Failed to cleanup upload: {e}")

        return jsonify({
            'success': True,
            'filename': output_zip.name,
            'download_url': f'/download/{output_zip.name}'
        })

    except zipfile.BadZipFile:
        return jsonify({'error': 'Invalid or corrupted ZIP file'}), 400
    except Exception as e:
        logger.error(f"Processing error: {e}", exc_info=True)
        return jsonify({'error': f'Processing failed: {str(e)}'}), 500

@app.route('/download/<path:filename>')
def download(filename):
    """Serve the processed ZIP file for download."""
    try:
        # Allow directory listing for manifest access
        if filename.endswith('/_manifest.md'):
            zip_name = filename.replace('/_manifest.md', '.zip')
            zip_path = OUTPUT_FOLDER / zip_name
            if not zip_path.exists():
                return 'File not found', 404

            # Extract manifest temporarily
            with zipfile.ZipFile(zip_path, 'r') as zf:
                try:
                    manifest_content = zf.read('_manifest.md').decode('utf-8')
                    return manifest_content, 200, {'Content-Type': 'text/markdown'}
                except KeyError:
                    return 'Manifest not found in archive', 404

        file_path = OUTPUT_FOLDER / filename
        if not file_path.exists():
            return 'File not found', 404

        return send_file(
            file_path,
            as_attachment=True,
            download_name=filename
        )
    except Exception as e:
        logger.error(f"Download error: {e}")
        return 'Internal server error', 500

@app.route('/health')
def health():
    """Health check endpoint."""
    return jsonify({'status': 'ok'})

# ============================================================================
# Main Entry Point
# ============================================================================

if __name__ == '__main__':
    print("""
================================================================================
  LLM Pre-Processing Pipeline Generator
================================================================================

  This application converts ZIP archives into LLM-optimized packages by
  applying the LLM Data-Type Hierarchy to each file.

  DEPENDENCIES:
  -------------
  System packages (install with):
      sudo apt-get update && sudo apt-get install -y pandoc tesseract-ocr \
          poppler-utils

  Python packages (install with):
      pip install Flask python-docx python-pptx openpyxl Pillow pytesseract

  OPTIONAL (for better PDF handling):
      pip install pdf2image

  USAGE:
  ------
      python app.py

  Then navigate to http://localhost:5000 in your browser.

================================================================================
""")

    app.run(host='0.0.0.0', port=5000, debug=False)
